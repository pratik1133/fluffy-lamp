"""
PPT Generator Script for Research Reports
==========================================
Version: 2.1.0 (2026-02-21) - Fixed current_para_text NameError
This script automates the generation of PowerPoint presentations
from Supabase data received via n8n webhook.

Template Structure (11 slides):
- Slide 1: Title (company_name, nse_symbol, bom_code, rating)
- Slide 2: Company Background
- Slide 3: Business Model
- Slide 4: Management Analysis
- Slide 5: Industry Overview
- Slide 6: Key Industry Tailwinds
- Slide 7: Demand Drivers
- Slide 8: Industry Risks
- Slide 9: Financials (summary_table text + 4 chart quadrants)
- Slide 10: Summary Charts (summary_charts text)
- Slide 11: Disclaimers
"""

import os
import re
import requests
from io import BytesIO
from datetime import datetime
from typing import Dict, Optional, Any, List, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE


class PPTGenerator:
    """
    A class to generate PowerPoint presentations from research report data.
    """

    # Chart image positions (slide index is 0-based)
    # Note: Slide 9 (index 8) has 4 financial chart quadrants AND the summary_table image
    # Slide 10 (index 9) has the chart_custom image
    CHART_POSITIONS = {
        'chart_profit_loss': {
            'slide': 8,  # Slide 9
            'position': {'left': 0.16, 'top': 1.4, 'width': 4.8, 'height': 2.2}
        },
        'chart_balance_sheet': {
            'slide': 8,  # Slide 9
            'position': {'left': 5.15, 'top': 1.4, 'width': 4.8, 'height': 2.2}
        },
        'chart_cash_flow': {
            'slide': 8,  # Slide 9
            'position': {'left': 0.10, 'top': 4.3, 'width': 4.8, 'height': 2.2}
        },
        'chart_ratio_analysis': {
            'slide': 8,  # Slide 9
            'position': {'left': 5.10, 'top': 4.3, 'width': 4.8, 'height': 2.2}
        },
        'summary_table': {
            'slide': 13,  # Slide 14 (Summary in Tables)
            'position': {'left': 0.5, 'top': 0.75, 'width': 9.0, 'height': 4.5}
        },
        'chart_custom': {
            'slide': 14,  # Slide 15 (Summary in Charts)
            'position': {'left': 0.5, 'top': 0.75, 'width': 9.0, 'height': 4.5}
        },
        'price_chart': {
             'slide': 13,  # Slide 14
             'position': {'left': 5.5, 'top': 1.6, 'width': 4.3, 'height': 2.5}
        },
    }

    def __init__(self, template_path: str):
        """Initialize the PPT Generator with a template."""
        self.template_path = template_path
        self.prs = None

    def load_template(self) -> None:
        """Load the PowerPoint template."""
        if not os.path.exists(self.template_path):
            raise FileNotFoundError(f"Template not found: {self.template_path}")
        self.prs = Presentation(self.template_path)
        print(f"  Loaded template with {len(self.prs.slides)} slides")

    def parse_markdown_to_text(self, markdown_text: str) -> str:
        """
        Convert markdown text to clean plain text.
        Preserves paragraph structure but removes markdown formatting.
        """
        if not markdown_text:
            return ""

        text = markdown_text

        # Remove markdown headers but keep the text
        text = re.sub(r'^#{1,6}\s*(.+)$', r'\1', text, flags=re.MULTILINE)

        # Convert bold/italic markers
        # We now PRESERVE bold markers (**) so they can be parsed by replace_shape_text for rich formatting
        # text = re.sub(r'\*\*\*(.+?)\*\*\*', r'\1', text)  # Bold italic
        # text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)      # Bold
        
        text = re.sub(r'(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)', r'\1', text)  # Italic (strip simple italic for now)
        text = re.sub(r'__(.+?)__', r'\1', text)          # Bold alt (strip)
        text = re.sub(r'(?<!_)_(?!_)(.+?)(?<!_)_(?!_)', r'\1', text)  # Italic alt

        # Remove link formatting but keep text
        text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)

        # Clean up excessive newlines (keep double newlines for paragraphs)
        text = re.sub(r'\n{3,}', '\n\n', text)
        
        # Remove leading/trailing whitespace from each line
        lines = [line.strip() for line in text.split('\n')]
        text = '\n'.join(lines)

        return text.strip()

    def download_image(self, url: str) -> Optional[BytesIO]:
        """Download an image from URL and return as BytesIO object."""
        if not url or url in ("[null]", "null", None, ""):
            return None

        try:
            print(f"    Downloading: {url[:60]}...")
            response = requests.get(url, timeout=30)
            response.raise_for_status()
            image_data = BytesIO(response.content)
            image_data.seek(0)
            return image_data
        except Exception as e:
            print(f"    Error downloading image: {e}")
            return None

    def find_shape_with_placeholder(self, placeholder_name: str):
        """
        Find the shape containing the placeholder text.
        Returns (slide, shape) tuple or (None, None) if not found.
        """
        placeholder_pattern = f"{{{{{placeholder_name}}}}}"
        
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                # Check full text of the shape
                full_text = ""
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        full_text += run.text
                
                if placeholder_pattern in full_text:
                    return slide, shape
        
        return None, None

    def replace_shape_text(self, shape, new_text: str, font_size: int = 12, bold: bool = False, align: str = None, color: Tuple[int, int, int] = None) -> bool:
        """
        Replace the entire text content of a shape with new text.
        Properly handles text frame formatting to prevent overflow and overlapping.
        """
        if not shape.has_text_frame:
            return False
        
        tf = shape.text_frame
        
        # Set text frame properties to prevent overflow
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # Auto-shrink text to fit
        
        # Set margins (Inches)
        tf.margin_top = Inches(0.25)
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_bottom = Inches(0.1)
        
        # Set vertical anchor
        if align and align.upper() == "CENTER":
             tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        else:
             tf.vertical_anchor = MSO_ANCHOR.TOP
        
        # Clear ALL existing paragraphs properly using XML manipulation
        from lxml import etree
        txBody = tf._txBody
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        all_paras = txBody.findall('a:p', nsmap)  # Direct children only
        print(f"    [DEBUG replace_shape_text] Existing paragraphs in shape: {len(all_paras)}")
        for p_elem in all_paras[1:]:
            txBody.remove(p_elem)
        if tf.paragraphs:
            tf.paragraphs[0].clear()
            # Remove ALL paragraph properties (indentation, margins, alignment, bullets) from first paragraph
            # so it starts fresh like newly-added paragraphs
            from pptx.oxml.ns import qn as qn_shape
            first_pPr = tf.paragraphs[0]._p.find(qn_shape('a:pPr'))
            if first_pPr is not None:
                tf.paragraphs[0]._p.remove(first_pPr)
        
        # Split new_text by newlines to create actual paragraphs
        lines = new_text.split('\n')
        print(f"    [DEBUG replace_shape_text] Text length: {len(new_text)} chars, splitting into {len(lines)} lines, font_size={font_size}")
        
        # For whole-shape text replacement, default to LEFT alignment for consistency
        effective_align = align if align else 'LEFT'
        
        # Use existing first paragraph for the first line
        if len(tf.paragraphs) > 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            
        if lines:
            self.replace_paragraph_with_markdown(p, lines[0], font_size, bold, effective_align, color)
            
        # Add additional paragraphs for subsequent lines
        for line in lines[1:]:
            p = tf.add_paragraph()
            self.replace_paragraph_with_markdown(p, line, font_size, bold, effective_align, color)

        # Verify final paragraph count
        final_paras = txBody.findall('a:p', nsmap)
        print(f"    [DEBUG replace_shape_text] Final paragraphs in shape: {len(final_paras)}")
        # Print shape dimensions
        try:
            w_inches = shape.width / 914400
            h_inches = shape.height / 914400
            print(f"    [DEBUG replace_shape_text] Shape size: {w_inches:.2f} x {h_inches:.2f} inches")
        except: pass

        return True

    def replace_paragraph_with_markdown(self, paragraph, text_content, font_size, bold, align, color=None):
        """
        Replaces paragraph text with Markdown-parsed runs.
        Supports both **bold** and *bold* syntax.
        Also applies heuristic to bold '- Label:' patterns.
        """
        paragraph.clear()
        
        # Ensure no bullet formatting from template
        from pptx.oxml.ns import qn
        pPr = paragraph._p.find(qn('a:pPr'))
        if pPr is not None:
            buNone = pPr.find(qn('a:buNone'))
            if buNone is None:
                for bu_elem in pPr.findall(qn('a:buChar')):
                    pPr.remove(bu_elem)
                for bu_elem in pPr.findall(qn('a:buAutoNum')):
                    pPr.remove(bu_elem)
                from lxml import etree
                etree.SubElement(pPr, qn('a:buNone'))
        
        # Auto-Bold Heuristic: Bold all "Label:" patterns in the text
        # This handles both "- Label: value" bullet patterns and "NSE:VALUE | BOM:VALUE" patterns
        # Find all "Word:" or "Multi Word:" patterns and wrap them in **
        # Use regex to find all label patterns (word(s) followed by colon)
        def bold_labels(text):
            # Match patterns like "Label:" at start or after "| " or "- " 
            # but skip if already inside ** markers
            result = re.sub(
                r'(?<!\*\*)(?:^|\| |(?<=\n))\s*[-•]?\s*([A-Za-z][A-Za-z &/\-\']+):',
                lambda m: m.group(0).replace(f"{m.group(1)}:", f"**{m.group(1)}:**"),
                text
            )
            return result
        
        if '**' not in text_content:
            # Only apply heuristic if there are no existing bold markers
            text_content = bold_labels(text_content)

        # Split key text by **...** OR *...* to identify bold sections
        # Regex matches **bold** OR *bold*
        parts = re.split(r'(\*\*.*?\*\*|\*.*?\*)', text_content)
        
        for part in parts:
            if not part:
                continue
                
            run = paragraph.add_run()
            
            # Check if this segment is wrapped in ** or *
            is_double_star = part.startswith('**') and part.endswith('**') and len(part) > 4
            is_single_star = part.startswith('*') and part.endswith('*') and len(part) > 2
            is_marked_bold = is_double_star or is_single_star
            
            # Strip markers if present
            if is_double_star:
                run_text = part[2:-2]
            elif is_single_star:
                run_text = part[1:-1]
            else:
                run_text = part
                
            run.text = run_text
            
            # Set font properties
            if font_size:
                 run.font.size = Pt(float(font_size)) 
            run.font.name = "Calibri"
            
            if bold:
                run.font.bold = True
            else:
                if is_marked_bold:
                    run.font.bold = True
                else:
                    run.font.bold = False 
            
            if color:
                run.font.color.rgb = RGBColor(*color)
            else:
                # Default to black text
                run.font.color.rgb = RGBColor(0, 0, 0)
                
        # Set alignment directly via XML (python-pptx may not write 'l' since LEFT is "default")
        if align:
            align_map = {
                'LEFT': 'l',
                'CENTER': 'ctr',
                'RIGHT': 'r',
                'JUSTIFY': 'just',
            }
            xml_align = align_map.get(align.upper(), 'l')
            # Get or create pPr element and set algn attribute directly
            from pptx.oxml.ns import qn as qn_align
            pPr_elem = paragraph._p.find(qn_align('a:pPr'))
            if pPr_elem is None:
                from lxml import etree
                pPr_elem = etree.SubElement(paragraph._p, qn_align('a:pPr'))
                paragraph._p.insert(0, pPr_elem)  # pPr must be first child
            pPr_elem.set('algn', xml_align)

    def replace_placeholder_with_image(self, placeholder_name: str, image_data: BytesIO) -> bool:
        """
        Find a shape containing {{placeholder_name}}, get its position/size,
        remove the placeholder text/shape, and insert the image in its place.
        """
        slide_obj, shape_obj = self.find_shape_with_placeholder(placeholder_name)
        
        if not slide_obj or not shape_obj:
            print(f"    -> Placeholder '{{{{{placeholder_name}}}}}' not found for image replacement.")
            return False
            
        # Get geometry
        left = shape_obj.left
        top = shape_obj.top
        width = shape_obj.width
        height = shape_obj.height
        
        # Insert image
        try:
            image_data.seek(0)
            slide_obj.shapes.add_picture(image_data, left, top, width, height)
            
            # Clear the placeholder text/shape so it doesn't show behind
            # We can't easily delete shapes in python-pptx without accessing xml, 
            # so we just clear the text.
            if shape_obj.has_text_frame:
                shape_obj.text_frame.clear()
                
            return True
        except Exception as e:
            print(f"    -> Error inserting image at placeholder: {e}")
            return False

    def find_and_replace_placeholder(self, placeholder_name: str, new_text: str, font_size: int = 12, bold: bool = False, align: str = None, color: Tuple[int, int, int] = None) -> int:
        """
        Find and replace {{placeholder_name}} with new text.
        Uses proper text replacement to avoid overflow issues.
        """
        placeholder_pattern = f"{{{{{placeholder_name}}}}}"
        replacements = 0

        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                tf = shape.text_frame
                
                # Get full text of the shape
                full_text = ""
                for para in tf.paragraphs:
                    full_text += ''.join(run.text for run in para.runs)
                
                # Check if placeholder exists
                if placeholder_pattern not in full_text:
                    continue
                
                # Robust check for "Whole Shape is Placeholder"
                # Remove whitespace from both and compare
                clean_full = "".join(full_text.split())
                clean_placeholder = "".join(placeholder_pattern.split())
                
                print(f"    [DEBUG] Placeholder '{placeholder_name}': full_text='{full_text[:80]}...', clean match={clean_full == clean_placeholder}, strip match={full_text.strip() == placeholder_pattern}")
                
                if clean_full == clean_placeholder or full_text.strip() == placeholder_pattern:
                    # This is a simple placeholder-only shape -> use replace_shape_text (multi-paragraph)
                    print(f"    [DEBUG] -> Taking WHOLE SHAPE branch (replace_shape_text)")
                    self.replace_shape_text(shape, new_text, font_size, bold, align, color)
                    replacements += 1
                else:
                    # Multiple placeholders or mixed content - do inline replacement
                    print(f"    [DEBUG] -> Taking INLINE branch (mixed content)")
                    for para in tf.paragraphs:
                        current_para_text = ''.join(run.text for run in para.runs)
                        if placeholder_pattern in current_para_text:
                            # Replace placeholder in the paragraph text, then re-render with markdown engine
                            # This ensures bold formatting is properly handled (**bold** = bold, rest = not bold)
                            new_para_text = current_para_text.replace(placeholder_pattern, new_text)
                            print(f"    [DEBUG] -> Inline replacing in para: '{current_para_text[:60]}...'")
                            self.replace_paragraph_with_markdown(para, new_para_text, font_size, bold, align, color)
                            replacements += 1

        return replacements

    def parse_markdown_table_to_data(self, markdown_text: str) -> List[List[str]]:
        """
        Parses a markdown table into a list of lists (rows of columns).
        Example input:
        | Header 1 | Header 2 |
        |---|---|
        | Row 1 Col 1 | Row 1 Col 2 |
        """
        if not markdown_text:
            return []
            
        lines = markdown_text.strip().split('\n')
        table_data = []
        
        for line in lines:
            # Skip separator lines (e.g. |---|---|)
            if '---' in line:
                continue
            # Skip empty lines
            if not line.strip():
                continue
            # Check if likely a row
            if '|' in line:
                # Split by pipe, strip whitespace
                row = [cell.strip() for cell in line.split('|')]
                # Filter out empty strings from leading/trailing pipes
                row = [cell for cell in row if cell]
                if row:
                    table_data.append(row)
                    
        return table_data

    def populate_table_shape(self, shape, data: List[List[str]], font_size: int = 12):
        """
        Populate a PowerPoint table shape with data (list of lists).
        """
        if not shape.has_table:
            return
            
        table = shape.table
        
        # Iterate over data rows
        for r_idx, row_data in enumerate(data):
            # If we run out of table rows, stop (or add rows if supported/needed)
            if r_idx >= len(table.rows):
                break
                
            for c_idx, cell_value in enumerate(row_data):
                # If we run out of table cols, stop
                if c_idx >= len(table.columns):
                    break
                    
                cell = table.cell(r_idx, c_idx)
                # Use same text replacement logic to handle formating
                self.replace_shape_text(cell, str(cell_value), font_size)

    def find_and_populate_table(self, placeholder_text: str, data: List[List[str]], font_size: int = 12) -> bool:
        """
        Find a table that contains the specific placeholder in its first cell (0,0)
        and populate it with the provided data.
        """
        target_pattern = f"{{{{{placeholder_text}}}}}"
        
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    # Check first cell for placeholder
                    try:
                        first_cell_text = shape.table.cell(0, 0).text_frame.text.strip()
                        if target_pattern in first_cell_text:
                            # Clear the placeholder logic from the first cell effectively by overwriting 
                            # when populating, or we treat the headers as row 0.
                            
                            print(f"    -> Found table with placeholder '{placeholder_text}' on Slide {self.prs.slides.index(slide)+1}")
                            self.populate_table_shape(shape, data, font_size)
                            return True
                    except Exception:
                        continue
        return False

    def create_table_on_slide(self, slide_idx: int, data: List[List[str]], 
                              left: float, top: float, width: float, height: float) -> bool:
        """
        Create a new table on a specific slide and populate it with styling.
        """
        if slide_idx >= len(self.prs.slides):
            return False
            
        if not data:
            return False
            
        rows = len(data)
        cols = len(data[0])
        
        try:
            slide = self.prs.slides[slide_idx]
            graphic_frame = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
            table = graphic_frame.table
            
            # Populate data and style
            for r, row_data in enumerate(data):
                for c, cell_value in enumerate(row_data):
                    cell = table.cell(r, c)
                    cell.text = str(cell_value)
                    
                    # Text Styling
                    if cell.text_frame.paragraphs:
                        para = cell.text_frame.paragraphs[0]
                        para.font.size = Pt(12)
                        if hasattr(para, 'font'):
                            para.font.name = "Calibri"
                        para.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                        
                    # Row Styling
                    cell.fill.solid()
                    if r == 0:
                        # Header: Dark Blue
                        cell.fill.fore_color.rgb = RGBColor(0, 51, 153)
                        if cell.text_frame.paragraphs:
                            para = cell.text_frame.paragraphs[0]
                            para.font.bold = True
                            para.font.color.rgb = RGBColor(255, 255, 255)
                    elif r % 2 == 0:
                        # Even index rows (2, 4, 6...) - Light Orange
                        cell.fill.fore_color.rgb = RGBColor(255, 235, 205)
                    else:
                        # Odd index rows (1, 3, 5...) - White
                        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            
            return True
        except Exception as e:
            print(f"    Error creating table: {e}")
            return False

    def add_image_to_slide(self, slide_idx: int, image_data: BytesIO,
                           left: float, top: float,
                           width: float, height: Optional[float] = None,
                           crop: Optional[Dict[str, float]] = None) -> bool:
        """Add an image to a specific slide with optional cropping."""
        if slide_idx >= len(self.prs.slides):
            print(f"    Warning: Slide {slide_idx + 1} does not exist")
            return False

        try:
            slide = self.prs.slides[slide_idx]
            image_data.seek(0)
            
            pic = None
            if height:
                pic = slide.shapes.add_picture(
                    image_data, 
                    Inches(left), Inches(top),
                    width=Inches(width), height=Inches(height)
                )
            else:
                pic = slide.shapes.add_picture(
                    image_data, 
                    Inches(left), Inches(top),
                    width=Inches(width)
                )
            
            # Apply cropping if provided
            if crop and pic:
                if 'left' in crop: pic.crop_left = crop['left']
                if 'top' in crop: pic.crop_top = crop['top']
                if 'right' in crop: pic.crop_right = crop['right']
                if 'bottom' in crop: pic.crop_bottom = crop['bottom']

            return True
        except Exception as e:
            print(f"    Error adding image: {e}")
            return False

    def add_debug_grid(self, slide_idx: int):
        """Add visual debug lines to the slide."""
        try:
            slide = self.prs.slides[slide_idx]
            from pptx.util import Inches
            from pptx.enum.shapes import MSO_CONNECTOR
            from pptx.dml.color import RGBColor

            # Draw Red Line at Top = 1.1 inches (Target top)
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(1.1), Inches(10), Inches(1.1)
            )
            line.line.color.rgb = RGBColor(255, 0, 0)
            line.line.width = Inches(0.05)

            # Draw Green Line at Top = 6.6 inches (Target bottom approx)
            line2 = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(6.6), Inches(10), Inches(6.6)
            )
            line2.line.color.rgb = RGBColor(0, 255, 0)
            line2.line.width = Inches(0.05)
            
            print(f"    DEBUG: Added red/green lines to Slide {slide_idx+1}")
            
            # Print slide dimensions
            print(f"    DEBUG: Slide width={self.prs.slide_width/914400} inches, height={self.prs.slide_height/914400} inches")

        except Exception as e:
            # Typically imports might fail if python-pptx version is old or structured differently
            # We try to import inside to be safe or just print error
            print(f"    Debug error: {e}")

    def calculate_font_size(self, text: str, max_chars: int = 2000) -> float:
        """
        Calculate appropriate font size based on text length.
        Longer text gets smaller font.
        """
        text_len = len(text)
        
        if text_len < 500:
            return 12.0  # Standard body text
        elif text_len < 1000:
            return 11.5
        elif text_len < 1500:
            return 11.0
        elif text_len < 2000:
            return 9.0
        elif text_len < 3000:
            return 8.0
        else:
            return 7.0

    def fetch_bom_code(self, symbol: str, company_name: str) -> str:
        """
        Attempt to fetch BSE numeric code using multiple sources:
        1. Hardcoded fallback list
        2. BSE India search API
        3. Screener.in
        4. Yahoo Finance (.BO symbol)
        """
        # Fallback list of common Indian stocks (NSE symbol -> BSE code)
        KNOWN_BSE_CODES = {
            'WIPRO': '507685', 'TCS': '532540', 'INFY': '500209',
            'RELIANCE': '500325', 'HDFCBANK': '500180', 'ICICIBANK': '532174',
            'SBIN': '500112', 'BHARTIARTL': '532454', 'ITC': '500875',
            'HINDUNILVR': '500696', 'KOTAKBANK': '500247', 'LT': '500510',
            'AXISBANK': '532215', 'ASIANPAINT': '500820', 'MARUTI': '532500',
            'TATAMOTORS': '500570', 'SUNPHARMA': '524715', 'TITAN': '500114',
            'BAJFINANCE': '500034', 'HCLTECH': '532281', 'BAJAJ-AUTO': '532977',
            'SWIGGY': '543842', 'ZOMATO': '543320',
        }
        
        symbol_upper = symbol.upper().strip() if symbol else ''
        
        # 1. Check hardcoded list
        if symbol_upper in KNOWN_BSE_CODES:
            print(f"    -> Found BSE code in fallback list: {KNOWN_BSE_CODES[symbol_upper]}")
            return KNOWN_BSE_CODES[symbol_upper]
        
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36'
        }
        
        # 2. Try BSE India search API
        try:
            for query in [symbol, company_name]:
                if not query:
                    continue
                print(f"    -> Searching BSE India for: {query}")
                bse_url = f"https://api.bseindia.com/BseIndiaAPI/api/ComHeadernewCompSearch/w?flag=suggestflag&scompany={query}"
                response = requests.get(bse_url, headers=headers, timeout=10)
                if response.status_code == 200:
                    results = response.json()
                    if results and len(results) > 0:
                        # BSE API returns list of dicts with 'scrip_cd' (BSE code)
                        for item in results:
                            scrip_cd = str(item.get('scrip_cd', item.get('SCRIP_CD', '')))
                            if scrip_cd and scrip_cd.isdigit():
                                print(f"    -> Found BSE code via BSE India: {scrip_cd}")
                                return scrip_cd
        except Exception as e:
            print(f"    -> BSE India search failed: {e}")
        
        # 3. Try Screener.in
        try:
            if symbol:
                print(f"    -> Searching Screener.in for: {symbol}")
                screener_url = f"https://www.screener.in/api/company/search/?q={symbol}"
                response = requests.get(screener_url, headers=headers, timeout=10)
                if response.status_code == 200:
                    results = response.json()
                    for item in results:
                        url = item.get('url', '')
                        # Screener URL format: /company/BSE_CODE/company-name/
                        # or /company/NSE_SYMBOL/
                        bse_id = str(item.get('bse_code', ''))
                        if bse_id and bse_id.isdigit():
                            print(f"    -> Found BSE code via Screener: {bse_id}")
                            return bse_id
        except Exception as e:
            print(f"    -> Screener search failed: {e}")
        
        # 4. Try Yahoo Finance (last resort)
        try:
            queries = [q for q in [symbol, company_name] if q]
            for query in queries:
                print(f"    -> Searching Yahoo Finance for: {query}")
                url = f"https://query2.finance.yahoo.com/v1/finance/search?q={query}&quotesCount=10&newsCount=0"
                response = requests.get(url, headers=headers, timeout=10)
                if response.status_code == 200:
                    data = response.json()
                    quotes = data.get('quotes', [])
                    for quote in quotes:
                        symbol_ticker = quote.get('symbol', '')
                        if symbol_ticker.endswith('.BO'):
                            bse_code = symbol_ticker.split('.')[0]
                            # Only return if it's numeric (actual BSE code)
                            if bse_code.isdigit():
                                print(f"    -> Found numeric BSE code via Yahoo: {bse_code}")
                                return bse_code
                            else:
                                print(f"    -> Yahoo returned non-numeric BSE ticker: {bse_code} (skipping)")
        except Exception as e:
            print(f"    -> Yahoo Finance search failed: {e}")
        
        print("    -> No numeric BSE code found from any source")
        return ' '

    def populate_from_data(self, data: Dict[str, Any]) -> Dict[str, bool]:
        """
        Populate the presentation with data from the Supabase record.
        """
        print("\n" + "=" * 60)
        print("POPULATING PRESENTATION")
        print("=" * 60)

        results = {}

        # --- PRE-CALCULATE METRICS (Growth, Margins) ---
        # Enrich data dictionary with calculated values so placeholders work
        years = [24, 25, 26, 27, 28]
        for y in years:
            prev_y = y - 1
            
            # Helper to parse
            def p_float(v):
                if v is None or v == '' or v == '-': return None
                try: 
                    if isinstance(v, str): v = v.replace(',', '')
                    return float(v)
                except: return None

            # 1. Margins
            rev = p_float(data.get(f'revenue_fy{y}'))
            ebit = p_float(data.get(f'ebitda_fy{y}'))
            if rev is not None and ebit is not None and rev != 0:
                margin = (ebit / rev) * 100
                data[f'ebitda_margin_fy{y}'] = "{:.1f}".format(margin)

            # 2. Growth
            for metric in ['revenue', 'ebitda', 'pat']:
                curr = p_float(data.get(f'{metric}_fy{y}'))
                prev = p_float(data.get(f'{metric}_fy{prev_y}'))
                
                if curr is not None and prev is not None and prev != 0:
                    growth = ((curr - prev) / abs(prev)) * 100
                    formatted_growth = "{:.1f}".format(growth)
                    data[f'{metric}_growth_fy{y}'] = formatted_growth
                    
                    # Alias 'revenue_growth' to 'sales_growth' for convenience
                    if metric == 'revenue':
                        data[f'sales_growth_fy{y}'] = formatted_growth

        # ===== TABLE POPULATION =====
        print("\n--- Table Population ---")
        
        # 1. Try to use explicit markdown if provided (highest priority)
        # 1. Try to use explicit markdown if provided (highest priority)
        financial_val = data.get('financial_performance', '')
        table_data = []
        financial_text_summary = ""
        
        has_markdown_table = False
        if financial_val and '|' in str(financial_val):
             has_markdown_table = True
             print("  Found markdown table in 'financial_performance'. Parsing...")
             table_data = self.parse_markdown_table_to_data(str(financial_val))
        else:
             # It acts as text summary if not a table
             if financial_val:
                 print("  'financial_performance' appears to be text summary.")
                 financial_text_summary = str(financial_val)
        
        # 2. If no markdown table found, try to construct from individual DB fields
        # 2. Financial Table Construction (New Logic - Slide 3)
        if not has_markdown_table:
            print("  Constructing table from new financial keys...")
            
            # Helper to safely parse float
            def safe_float(val):
                if val is None or val == '' or val == '-': return None
                try:
                    if isinstance(val, str):
                        val = val.replace(',', '').strip()
                    return float(val)
                except:
                    return None

            # Helper to calculate YoY Growth %
            def calc_growth(current_val, prev_val):
                curr = safe_float(current_val)
                prev = safe_float(prev_val)
                if curr is not None and prev is not None and prev != 0:
                    growth = ((curr - prev) / abs(prev)) * 100
                    return "{:.1f}".format(growth)
                return "-"

            # Helper to calculate Margin %
            def calc_margin(numerator_val, denominator_val):
                num = safe_float(numerator_val)
                den = safe_float(denominator_val)
                if num is not None and den is not None and den != 0:
                    margin = (num / den) * 100
                    return "{:.1f}".format(margin)
                return "-"

            # Helper to safely get numeric value formatted (direct lookup)
            def get_val_fmt(key, fmt="{:,.0f}", multiplier=1.0):
                val = safe_float(data.get(key))
                if val is not None:
                     return fmt.format(val * multiplier)
                return "-"

            # Base Metrics
            sales_24 = data.get('revenue_fy24')
            sales_25 = data.get('revenue_fy25')
            sales_26 = data.get('revenue_fy26')
            sales_27 = data.get('revenue_fy27')
            sales_28 = data.get('revenue_fy28')

            ebitda_24 = data.get('ebitda_fy24')
            ebitda_25 = data.get('ebitda_fy25')
            ebitda_26 = data.get('ebitda_fy26')
            ebitda_27 = data.get('ebitda_fy27') 
            ebitda_28 = data.get('ebitda_fy28')

            pat_24 = data.get('pat_fy24')
            pat_25 = data.get('pat_fy25')
            pat_26 = data.get('pat_fy26')
            pat_27 = data.get('pat_fy27')
            pat_28 = data.get('pat_fy28')

            # Headers
            headers = ["Particulars", "FY24A", "FY25A", "FY26E", "FY27E", "FY28E"]
            
            # Data Rows with Calculations
            rows = [
                # Sales (Direct)
                ["Sales", 
                 get_val_fmt('revenue_fy24'), get_val_fmt('revenue_fy25'), 
                 get_val_fmt('revenue_fy26'), get_val_fmt('revenue_fy27'), get_val_fmt('revenue_fy28')],
                
                # Sales Growth (Calculated)
                ["YoY% growth", 
                 calc_growth(sales_24, data.get('revenue_fy23')), # Need FY23 for FY24 growth, else use provided key
                 calc_growth(sales_25, sales_24), 
                 calc_growth(sales_26, sales_25), 
                 calc_growth(sales_27, sales_26), 
                 calc_growth(sales_28, sales_27)],

                # EBITDA (Direct)
                ["EBITDA", 
                 get_val_fmt('ebitda_fy24'), get_val_fmt('ebitda_fy25'), 
                 get_val_fmt('ebitda_fy26'), get_val_fmt('ebitda_fy27'), get_val_fmt('ebitda_fy28')],

                # EBITDA Margin (Calculated: EBITDA / Sales)
                ["% Margin", 
                 calc_margin(ebitda_24, sales_24), calc_margin(ebitda_25, sales_25), 
                 calc_margin(ebitda_26, sales_26), calc_margin(ebitda_27, sales_27), calc_margin(ebitda_28, sales_28)],

                # EBITDA Growth (Calculated)
                ["YoY% growth", 
                 calc_growth(ebitda_24, data.get('ebitda_fy23')),
                 calc_growth(ebitda_25, ebitda_24), 
                 calc_growth(ebitda_26, ebitda_25), 
                 calc_growth(ebitda_27, ebitda_26), 
                 calc_growth(ebitda_28, ebitda_27)],

                # PAT (Direct)
                ["PAT", 
                 get_val_fmt('pat_fy24'), get_val_fmt('pat_fy25'), 
                 get_val_fmt('pat_fy26'), get_val_fmt('pat_fy27'), get_val_fmt('pat_fy28')],

                # PAT Growth (Calculated)
                ["YoY% growth", 
                 calc_growth(pat_24, data.get('pat_fy23')),
                 calc_growth(pat_25, pat_24), 
                 calc_growth(pat_26, pat_25), 
                 calc_growth(pat_27, pat_26), 
                 calc_growth(pat_28, pat_27)],
                
                # P/E (Direct - difficult to calculate without Price history)
                ["P/E", 
                 get_val_fmt('pe_fy24', "{:.1f}"), get_val_fmt('pe_fy25', "{:.1f}"), 
                 get_val_fmt('pe_fy26', "{:.1f}"), get_val_fmt('pe_fy27', "{:.1f}"), get_val_fmt('pe_fy28', "{:.1f}")],

                # P/B (Direct)
                ["P/B", 
                 get_val_fmt('pb_fy24', "{:.1f}"), get_val_fmt('pb_fy25', "{:.1f}"), 
                 get_val_fmt('pb_fy26', "{:.1f}"), get_val_fmt('pb_fy27', "{:.1f}"), get_val_fmt('pb_fy28', "{:.1f}")] 
            ]
            
            table_data = [headers] + rows
             
            # Create Table on Slide 3 (Index 2)
            # Position: Below chart. Chart Bottom = ~3.2. Start Table at 3.3.
            # Width matches chart: 4.83.
            # DISABLED: User wants to use manual placeholders + custom table in template
            # if self.create_table_on_slide(2, table_data, left=5.0, top=3.3, width=4.83, height=2.5):
            #     print(f"  Financial Table (Slide 3): [OK] Created")
            # else:
            #     print(f"  Financial Table (Slide 3): [FAILED] Could not create")

        # OLD LOGIC DISABLED (Skipped)
        if False and ('revenue_fy2024' in data or 'revenue_ttm' in data):
            print("  Constructing table from equity_universe fields...")
            
            # Helper to safely get numeric value formatted
            def get_val(key, fmt="{:,.0f}"):
                val = data.get(key)
                if val is None: return "-"
                try:
                    return fmt.format(float(val))
                except:
                    return str(val)

            # Define the structure based on your screenshot
            # Header Row
            headers = ["Particulars", "FY24A", "FY25A", "FY26E", "FY27E", "FY28E"]
            
            # Data Rows
            rows = [
                # Sales
                ["Sales", 
                 get_val('revenue_fy2024'), get_val('revenue_fy2025'), 
                 get_val('revenue_fy2026e'), get_val('revenue_fy2027e'), get_val('revenue_fy2028e')],
                
                # Sales Growth (YoY %) - You might need to calculate this if not in DB
                ["YoY% growth", 
                 get_val('sales_growth_yoy_qtr', "{:.1f}"), "-", "-", "-", "-"], 

                # EBITDA
                ["EBITDA", 
                 get_val('ebitda_fy2024'), get_val('ebitda_fy2025'), 
                 get_val('ebitda_fy2026e'), get_val('ebitda_fy2027e'), get_val('ebitda_fy2028e')],

                # EBITDA Margin (%)
                ["% Margin", 
                 get_val('ebitda_margin_fy2024', "{:.1f}"), get_val('ebitda_margin_fy2025', "{:.1f}"), 
                 get_val('ebitda_margin_fy2026e', "{:.1f}"), get_val('ebitda_margin_fy2027e', "{:.1f}"), get_val('ebitda_margin_fy2028e', "{:.1f}")],

                # PAT
                ["PAT", 
                 get_val('pat_fy2024'), get_val('pat_fy2025'), 
                 get_val('pat_fy2026e'), get_val('pat_fy2027e'), get_val('pat_fy2028e')],

                # PAT Growth
                ["YoY% growth", 
                 get_val('pat_growth_qoq', "{:.1f}"), "-", "-", "-", "-"], # Using QoQ as placeholder if YoY missing

                # P/E
                ["P/E", 
                 get_val('pe_ttm', "{:.1f}"), get_val('pe_fy2025', "{:.1f}"), # Note: pe_fy2025 might not exist, check DB keys
                 get_val('pe_fy2026e', "{:.1f}"), get_val('pe_fy2027e', "{:.1f}"), get_val('pe_fy2028e', "{:.1f}")],

                # P/B (Book Value) - We have book_value, need P/B calculation or field
                ["P/B", 
                 "-", "-", "-", "-", "-"] 
            ]
            
            table_data = [headers] + rows
        
        # 3. Populate if we have data
        # DISABLED: User wants to replace {{financial_table}} with an image, not populate a table
        # if table_data:
        #     success = self.find_and_populate_table('financial_table', table_data, font_size=10)
        #     print(f"  Financial Table: {'[OK] Populated' if success else '[FAILED] Table placeholder {{financial_table}} not found'}")
        # else:
        #     print("  Financial Table: No data found (markdown or DB fields)")
        
        # ===== TEXT REPLACEMENTS =====
        print("\n--- Text Replacements ---")
        
        # Get or fetch BOM code (must be numeric like "507685")
        bom_code = data.get('bom_code', '')
        # Check if bom_code is valid (should be numeric)
        is_valid_bom = bom_code and str(bom_code).strip().isdigit()
        
        if not is_valid_bom:
            print(f"  BOM Code '{bom_code}' is invalid (not numeric). Fetching from Yahoo Finance...")
            symbol = data.get('nse_symbol', data.get('symbol', ''))
            name = data.get('company_name', '')
            bom_code = self.fetch_bom_code(symbol, name)
            print(f"  -> Found: {bom_code}" if bom_code.strip() else "  -> Not found")
        else:
            print(f"  BOM Code: {bom_code} (provided)")
        
        # Get rating, default to N/A if missing
        rating = data.get('rating', '')
        if not rating or str(rating).strip() == '':
            rating = 'N/A'
        print(f"  DEBUG: Rating/Recommendation value: '{rating}'")
        
        # Define placeholder mappings with their data sources
        text_mappings = [
            # === SLIDE 1: Title Slide ===
            ('company_name', data.get('company_name', ''), 40, {'bold': True, 'align': 'CENTER', 'color': (255, 255, 255)}),
            ('nse_symbol', data.get('nse_symbol', data.get('symbol', '')), 14, {'align': 'CENTER'}),
            ('bom_code', bom_code, 14, {'align': 'CENTER'}),
            ('recommendation', rating, 14, {'align': 'CENTER'}),
            ('today_date', data.get('today_date', datetime.now().strftime('%Y-%m-%d')), 14, {'align': 'CENTER'}),

            # === SLIDE 2: Critical Summary (cs_*) + Masterheading ===
            # DB: cs_masterheading -> Template: {{Masterheading_h}} (if it exists in template)
            ('Masterheading_h', data.get('cs_masterheading') or data.get('masterheading_h') or "Company Insider", 20, {'bold': True, 'align': 'CENTER', 'color': (255, 255, 255)}),
            # DB: cs_marketing_positioning -> Template: {{cs_market_positioning}}
            ('cs_market_positioning', self.parse_markdown_to_text(data.get('cs_marketing_positioning', data.get('market_positioning', ''))), 10),
            # DB: cs_financial_performance -> Template: {{cs_financial_performance}}
            ('cs_financial_performance', self.parse_markdown_to_text(data.get('cs_financial_performance', financial_text_summary)), 10),
            # DB: cs_grow_outlook -> Template: {{cs_grow_outlook}}
            ('cs_grow_outlook', self.parse_markdown_to_text(data.get('cs_grow_outlook', data.get('growth_outlook', ''))), 10),
            # DB: cs_value_and_recommendation -> Template: {{cs_valuation_recommendation}}
            ('cs_valuation_recommendation', self.parse_markdown_to_text(data.get('cs_value_and_recommendation', data.get('valuation_recommendation', ''))), 10),
            # DB: cs_key_risks -> Template: {{cs_key_risks}}
            ('cs_key_risks', self.parse_markdown_to_text(data.get('cs_key_risks', data.get('key_risks', ''))), 10),

            # === SLIDE 3: Company Background ===
            ('Company_Background_h', data.get('company_background_h') or "Company Background", 20, {'bold': True, 'align': 'CENTER', 'color': (255, 255, 255)}),
            ('company_background', self.parse_markdown_to_text(data.get('company_background', '')), 11),

            # === SLIDE 4: Business Model ===
            ('Business_Model_Explanation_h', data.get('business_model_h') or "Business Model", 20, {'bold': True, 'align': 'CENTER', 'color': (255, 255, 255)}),
            ('business_model', self.parse_markdown_to_text(data.get('business_model', '')), 11),

            # === SLIDE 5: Management Analysis ===
            ('Management_Analysis_h', data.get('management_analysis_h') or "Management Analysis", 20, {'bold': True, 'align': 'CENTER', 'color': (255, 255, 255)}),
            ('management_analysis', self.parse_markdown_to_text(data.get('management_analysis', '')), 11),

            # === SLIDE 6: Industry Overview ===
            ('Industry_Overview_h', data.get('industry_overview_h') or "Industry Overview", 20, {'bold': True, 'align': 'CENTER', 'color': (255, 255, 255)}),
            ('industry_overview', self.parse_markdown_to_text(data.get('industry_overview', '')), 11),

            # === SLIDE 7: Key Industry Tailwinds ===
            ('Key_Industry_Tailwinds_h', data.get('industry_tailwinds_h') or "Key Industry Tailwinds", 20, {'bold': True, 'align': 'CENTER', 'color': (255, 255, 255)}),
            ('industry_tailwinds', self.parse_markdown_to_text(data.get('industry_tailwinds', data.get('key_industry', ''))), 11),

            # === SLIDE 8: Demand Drivers ===
            ('Demand_drivers_h', data.get('demand_drivers_h') or "Demand Drivers", 20, {'bold': True, 'align': 'CENTER', 'color': (255, 255, 255)}),
            ('demand_drivers', self.parse_markdown_to_text(data.get('demand_drivers', '')), 11),

            # === SLIDE 9: Industry Risks ===
            ('Industry_Risks_h', data.get('industry_risks_h') or "Industry Risks", 20, {'bold': True, 'align': 'CENTER', 'color': (255, 255, 255)}),
            ('industry_risk', self.parse_markdown_to_text(data.get('industry_risks', data.get('industry_risk', ''))), 11),

            # === Extra fields (not in template but mapped for future use) ===
            ('market_positioning', self.parse_markdown_to_text(data.get('market_positioning', '')), 11),
            ('financial_performance', financial_text_summary, 11),
            ('grow_outlook', self.parse_markdown_to_text(data.get('growth_outlook', '')), 11),
            ('valuation_recommendation', self.parse_markdown_to_text(data.get('valuation_recommendation', '')), 11),
            ('key_risks', self.parse_markdown_to_text(data.get('key_risks', '')), 11),
            ('company_insider', self.parse_markdown_to_text(data.get('company_insider', '')), 11),
            ('cs_company_insider', self.parse_markdown_to_text(data.get('cs_company_insider', data.get('company_insider', ''))), 10),

            # === New Sections ===
            ('cs_investment_rationale', self.parse_markdown_to_text(data.get('cs_investment_rationale', '')), 10),
            ('cs_corporate_governance', self.parse_markdown_to_text(data.get('cs_corporate_governance', '')), 10),
            ('cs_saarthi_framework', self.parse_markdown_to_text(data.get('cs_saarthi_framework', '')), 10),
            ('cs_entry_review_exit_strategy', self.parse_markdown_to_text(data.get('cs_entry_review_exit_strategy', '')), 10),
            ('cs_scenario_analysis', self.parse_markdown_to_text(data.get('cs_scenario_analysis', '')), 10),

            # === New Metrics ===
            ('rating', str(data.get('cs_rating', '')), 12),
            ('target_price', str(data.get('cs_target_price', '')), 12),
            ('upside_percentage', str(data.get('cs_upside_percentage', '')), 12),
            ('market_cap', str(data.get('cs_market_cap', '')), 12),
            ('market_cap_category', str(data.get('cs_market_cap_category', '')), 12),
            ('current_market_price', str(data.get('cs_current_market_price', '')), 12),

            # === Scripts ===
            ('podcast_script', self.parse_markdown_to_text(data.get('podcast_script', '')), 11),
            ('video_script', self.parse_markdown_to_text(data.get('video_script', '')), 11),
            
            # === Image placeholders cleared (images inserted by fixed positioning) ===
            ('prize_chart', ' ', None),
            ('financial_table', ' ', None),
            ('summary_table', ' ', None),
            ('chart_custom', ' ', None),
        ]
        
        # --- DYNAMIC FINANCIAL PLACEHOLDERS ---
        # Allows user to use {{revenue_fy24}}, {{ebitda_margin_fy25}}, etc. in PPT if they wish
        financial_prefixes = ['revenue', 'sales', 'ebitda', 'pat', 'pe', 'pb']
        years = ['fy24', 'fy25', 'fy26', 'fy27', 'fy28', 'ttm']
        suffixes = ['', '_growth', '_margin']
        
        for p in financial_prefixes:
            for y in years:
                for s in suffixes:
                    key = f"{p}{s}_{y}" # e.g. revenue_fy24, ebitda_margin_fy25
                    if key in data:
                        # Add to mappings if data exists
                        val = str(data.get(key, '-'))
                        # Format numbers nicely if possible
                        try:
                            fval = float(val.replace(',', ''))
                            # If margin or small number, 1 decimal. If large, 0 decimals.
                            if 'margin' in s or 'growth' in s or fval < 100:
                                val = "{:.1f}".format(fval)
                            else:
                                val = "{:,.0f}".format(fval)
                        except:
                            pass
                        text_mappings.append((key, val, 12)) # Font size 12 for "small placeholders"

        for item in text_mappings:
            placeholder = item[0]
            value = item[1]
            fixed_font_size = item[2]
            formatting = item[3] if len(item) > 3 else {}
            
            if value:
                if fixed_font_size:
                    font_size = fixed_font_size
                else:
                    font_size = self.calculate_font_size(value)
                
                bold = formatting.get('bold', False)
                align = formatting.get('align', None)
                color = formatting.get('color', None)
                
                count = self.find_and_replace_placeholder(placeholder, value, font_size, bold, align, color)
                results[placeholder] = count > 0
                char_info = f"{len(value)} chars, {font_size}pt"
                status = f"[OK] Replaced ({char_info})" if count > 0 else "[MISSING] Placeholder not found"
                print(f"  {placeholder}: {status}")
            else:
                pass

        # ===== IMAGE INSERTIONS =====
        print("\n--- Image Insertions ---")
        
        # 1. Dynamic Replacement (using placeholders)
        # Added financial_table here to replace {{financial_table}} with image from Supabase
        dynamic_images = {
             # 'financial_table': {
             #     'url': data.get('financial_table'),
             #     'placeholder': 'financial_table'
             # }
        }
        
        for name, info in dynamic_images.items():
            url = info['url']
            placeholder = info['placeholder']
            
            if url and url not in ("[null]", "null", None, ""):
                print(f"  {name} (via {{{{{placeholder}}}}}):")
                image_data = self.download_image(url)
                if image_data:
                    success = self.replace_placeholder_with_image(placeholder, image_data)
                    results[name] = success
                    print(f"    -> {'[OK] Replaced placeholder' if success else '[FAILED] Placeholder not found'}")
                else:
                     results[name] = False
                     print("    -> [FAILED] Download failed")
            else:
                results[name] = False
                print(f"  {name}: [MISSING] No URL provided")

        # 2. Fixed Position Replacement
        fixed_images = {
            'chart_custom': { 
                'url': data.get('chart_custom'), 
                'slide': 14, # Slide 15 (Index 14)
                'pos': {'left': 0.5, 'top': 0.75, 'width': 9.0, 'height': 4.5} 
            },
            'price_chart_slide2': { 
                'url': data.get('price_chart'), 
                'slide': 1, # Slide 2 (Index 1)
                # User provided Size (cm->inch): W=12.27->4.83, H=5.08->2.0
                # Position estimated (Top Right)
                'pos': {'left': 5.0, 'top': 0.75, 'width': 4.83, 'height': 2.0} 
            },
            'financial_table_slide2': {
                'url': data.get('financial_table'),
                'slide': 1, # Slide 2 (Index 1)
                # User requested Width 12cm -> 4.72"
                'pos': {'left': 5.07, 'top': 2.81, 'width': 4.72, 'height': 2.06}
            },
            'summary_table_slide10': { 
                'url': data.get('summary_table'), 
                'slide': 13, # Slide 14 (Index 13)
                # User provided crop pos (cm->inch, calculated):
                # Left=1.52cm -> 0.60"
                # Top=1.73cm -> 0.68"
                # Width=21.18cm -> 8.34"
                # Height=12.2cm -> 4.80"
                'pos': {'left': 0.60, 'top': 0.68, 'width': 8.34, 'height': 4.80},
                # Calculated crop percentages (Offsets considered):
                # W_total=21.74, W_crop=21.18, OffX=0.36 -> R_crop=3.0%, L_crop=0%
                # H_total=15.69, H_crop=12.2, OffY=-0.34 -> T_crop=13.3%, B_crop=8.9%
                'crop': {'left': 0.0, 'right': 0.030, 'top': 0.133, 'bottom': 0.089}
            },
        }

        for name, info in fixed_images.items():
            url = info['url']
            if url and url not in ("[null]", "null", None, ""):
                print(f"  {name}:")
                image_data = self.download_image(url)
                if image_data:
                    slide_idx = info['slide']
                    pos = info['pos']
                    if pos:
                        crop = info.get('crop')
                        success = self.add_image_to_slide(
                            slide_idx, image_data, 
                            left=pos.get('left'), top=pos.get('top'), 
                            width=pos.get('width'), height=pos.get('height'),
                            crop=crop
                        )
                        results[name] = success
                        print(f"    -> Slide {slide_idx+1}: {'[OK] Added' if success else '[FAILED]'}")
                else:
                    results[name] = False
                    print(f"    -> [FAILED] Download failed")
            else:
                results[name] = False
                print(f"  {name}: [MISSING] No URL provided")

        return results

    def save(self, output_path: str) -> str:
        """Save the presentation to a file."""
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        self.prs.save(output_path)
        print(f"\n[OK] Presentation saved to: {output_path}")
        return output_path


def generate_report_ppt(data: Dict[str, Any], 
                        template_path: str,
                        output_dir: str = "./output") -> str:
    """
    Main function to generate a PowerPoint report.
    """
    # Generate output filename
    report_id = data.get('report_id', 'unknown')
    symbol = data.get('symbol', data.get('nse_symbol', 'report'))
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    
    # Clean report_id for filename
    report_id_clean = report_id[:8] if len(report_id) > 8 else report_id
    output_filename = f"{symbol}_{report_id_clean}_{timestamp}.pptx"
    output_path = os.path.join(output_dir, output_filename)

    # Create generator and process
    generator = PPTGenerator(template_path)
    generator.load_template()
    results = generator.populate_from_data(data)
    generator.save(output_path)

    # Summary
    successful = sum(1 for v in results.values() if v)
    total = len(results)
    print(f"\n{'=' * 60}")
    print(f"SUMMARY: {successful}/{total} fields processed successfully")
    print(f"{'=' * 60}")

    return output_path


# ============================================================
# EXAMPLE USAGE AND TESTING
# ============================================================
if __name__ == "__main__":
    # Example data structure (as received from n8n/Supabase)
    example_data = {
        "report_id": "c49b2aa1-80eb-4436-b14c-2a74d7966feb",
        "company_name": "Vedanta Ltd.",
        "nse_symbol": "VEDL",
        "bom_code": "500295",
        "rating": "BUY",
        "today_date": "2026-02-07",
        "company_background": """Company Background

Vedanta Limited is a globally diversified natural resources company with business operations in India, South Africa, Namibia, and Australia. The company is one of the world's largest diversified natural resources companies.

Key Business Segments:
• Zinc Business: One of the largest integrated producers of zinc-lead
• Aluminum Business: India's largest aluminum producer
• Oil and Gas: Significant crude oil producer in India
• Iron Ore: Mining operations in Goa and Karnataka
• Copper: Copper smelting and refining operations

History and Evolution:
Founded in 1976, Vedanta has grown through strategic acquisitions and organic expansion. The company was originally focused on mining and has diversified into various natural resources over the decades.

Market Position:
Vedanta holds leadership positions in multiple segments of the Indian natural resources industry, with significant global presence in key commodities.""",

        "business_model": """Business Model Explanation

Revenue Streams:
Vedanta generates revenue through multiple integrated business segments including mining operations, smelting and refining, and oil and gas production.

1. Mining Operations
• Extraction of zinc, lead, silver, iron ore
• Open-pit and underground mining operations
• Mineral processing and concentration

2. Smelting and Refining
• Aluminum smelting operations
• Copper cathode production
• Zinc and lead refining

3. Oil and Gas Production
• Crude oil extraction from Rajasthan fields
• Natural gas production

Value Chain Integration:
The company maintains vertical integration across exploration, mining, processing, and marketing. This integration provides cost advantages and supply chain control.

Key Competitive Advantages:
• Low-cost production capabilities
• Diverse commodity portfolio reducing risk
• Strong operational expertise
• Strategic asset locations""",

        "management_analysis": """Management Analysis

Leadership Team:
Anil Agarwal - Chairman: Founder and visionary leader with over 40 years of industry experience, known for bold strategic decisions.

Key Management Metrics:
• Experience: Excellent
• Track Record: Strong
• Corporate Governance: Good
• Capital Allocation: Above Average

Strategic Direction:
The management has outlined a clear growth strategy focusing on capacity expansion in aluminum and zinc, exploration and development of new resources, ESG improvements and sustainability initiatives, and digital transformation of operations.""",

        "industry_overview": """Industry Overview

Industry Size & Structure:
The mining and metals industry is a significant contributor to the global economy. The Total Addressable Market for this sector is vast, driven by demand for essential metals such as aluminum, copper, zinc, and iron ore.

Market Dynamics:
• Total global mining market: $2.1 trillion
• Base metals segment: $650 billion
• Expected CAGR: 4.5% (2024-2030)

Indian Market Position:
• India is the 3rd largest producer of coal
• 4th largest producer of iron ore
• Significant growth potential in base metals""",

        # --- MOCK FINANCIAL DATA (equity_universe fields) ---
        "revenue_fy2024": 150000, "revenue_fy2025": 165000, "revenue_fy2026e": 180000, "revenue_fy2027e": 200000, "revenue_fy2028e": 225000,
        "sales_growth_yoy_qtr": 12.5,
        "ebitda_fy2024": 45000, "ebitda_fy2025": 50000, "ebitda_fy2026e": 55000, "ebitda_fy2027e": 62000, "ebitda_fy2028e": 70000,
        "ebitda_margin_fy2024": 30.0, "ebitda_margin_fy2025": 30.3, "ebitda_margin_fy2026e": 30.5, "ebitda_margin_fy2027e": 31.0, "ebitda_margin_fy2028e": 31.1,
        "pat_fy2024": 12000, "pat_fy2025": 14000, "pat_fy2026e": 16000, "pat_fy2027e": 19000, "pat_fy2028e": 23000,
        "pat_growth_qoq": 15.2,
        "pe_ttm": 15.4, "pe_fy2025": 14.2, "pe_fy2026e": 12.5, "pe_fy2027e": 10.8, "pe_fy2028e": 9.2,

        "industry_tailwinds": """Key Industry Tailwinds

Structural Growth Drivers:

1. Infrastructure Development
• Government's infrastructure push (PM Gati Shakti)
• National Infrastructure Pipeline: ₹111 lakh crore investment
• Increased demand for steel, aluminum, and copper

2. Electric Vehicle Revolution
• EV adoption driving copper and aluminum demand
• Battery metals gaining importance
• India's EV sales growing at 40%+ CAGR

3. Renewable Energy Expansion
• Solar and wind capacity additions
• Transmission infrastructure build-out
• Energy storage requirements

4. Manufacturing Renaissance
• PLI schemes attracting investment
• China+1 strategy benefiting India

5. Urbanization Trends
• 40% urbanization currently, growing to 50% by 2030
• Housing and construction demand

Government Policy Support:
National Mineral Policy 2019, mining reforms and auction regime, export restrictions protecting domestic supply.""",

        "demand_drivers": """Demand Drivers for Vedanta Ltd.

End-User Industries:

1. Construction & Infrastructure (35% of demand)
• Real estate development
• Road and highway construction
• Port and airport development

2. Automotive Sector (20% of demand)
• Passenger and commercial vehicles
• Two-wheelers and EV components

3. Electrical & Electronics (18% of demand)
• Power cables and wiring
• Consumer electronics

4. Packaging Industry (12% of demand)
• Beverage cans
• Food and pharmaceutical packaging

5. Other Industries (15% of demand)
• Aerospace and defense
• Industrial machinery

Growth Outlook by Segment:
• Construction: 8% current, 10% outlook
• Automotive: 12% current, 15% outlook
• Electronics: 15% current, 18% outlook""",

        "industry_risks": """Industry Risks

Regulatory & Policy Risks:
• Environmental Regulations: Stricter emission norms, water usage restrictions
• Government Policy Changes: Export duty variations, royalty rate changes
• Impact: HIGH | Likelihood: MEDIUM

Market Risks:
• Commodity Price Volatility: Global demand-supply dynamics, currency fluctuations
• Competition Intensity: New capacity additions, import competition
• Impact: HIGH | Likelihood: HIGH

Operational Risks:
• Resource Depletion: Mine life limitations, grade deterioration
• Labor and Social Issues: Union negotiations, community relations
• Impact: MEDIUM | Likelihood: MEDIUM

Mitigation Strategies:
• Diversified commodity portfolio
• Long-term contracts with customers
• Hedging strategies for currency and commodities
• Strong community engagement programs""",

        "summary_table": "https://bmpvcjbfeyvkkbvclwkb.supabase.co/storage/v1/object/public/charts/VEDL/summary_table_example.png",
        "chart_custom": "https://bmpvcjbfeyvkkbvclwkb.supabase.co/storage/v1/object/public/charts/VEDL/custom_chart_example.png",
        "chart_profit_loss": "https://bmpvcjbfeyvkkbvclwkb.supabase.co/storage/v1/object/public/charts/VEDL/profit_loss_20260207_111704.png",
        "chart_balance_sheet": "https://bmpvcjbfeyvkkbvclwkb.supabase.co/storage/v1/object/public/charts/VEDL/balance_sheet_20260207_111704.png",
        "chart_cash_flow": "https://bmpvcjbfeyvkkbvclwkb.supabase.co/storage/v1/object/public/charts/VEDL/cash_flow_20260207_111705.png",
        "chart_ratio_analysis": "https://bmpvcjbfeyvkkbvclwkb.supabase.co/storage/v1/object/public/charts/VEDL/ratios_20260207_111705.png",
        "chart_summary": "https://bmpvcjbfeyvkkbvclwkb.supabase.co/storage/v1/object/public/charts/VEDL/custom_chart_20260207_111702.png",
    }

    # Get the directory of this script
    script_dir = os.path.dirname(os.path.abspath(__file__))
    template_path = os.path.join(script_dir, "master_template.pptx")
    output_dir = os.path.join(script_dir, "output")

    print("=" * 60)
    print("PPT GENERATOR - Research Report Automation")
    print("=" * 60)
    print(f"\nTemplate: {template_path}")
    print(f"Output Directory: {output_dir}")

    try:
        output_file = generate_report_ppt(
            data=example_data,
            template_path=template_path,
            output_dir=output_dir
        )
        print(f"\n{'=' * 60}")
        print(f"SUCCESS! Report generated: {output_file}")
        print("=" * 60)
    except Exception as e:
        print(f"\nERROR: {e}")
        import traceback
        traceback.print_exc()